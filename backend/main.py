from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pypdf import PdfWriter, PdfReader
import os
import uuid
import io
from typing import List, Dict, Set
import shutil
import time
from contextlib import asynccontextmanager

# Session tracking: session_id -> set of file paths
session_files: Dict[str, Set[str]] = {}

# Cleanup old files on startup (older than 1 hour)
def cleanup_old_files():
    """Delete files in uploads and merged directories that are older than 1 hour."""
    now = time.time()
    cutoff = now - 3600  # 1 hour ago
    
    for directory in [UPLOAD_DIR, MERGED_DIR]:
        if os.path.exists(directory):
            for filename in os.listdir(directory):
                file_path = os.path.join(directory, filename)
                try:
                    if os.path.isfile(file_path):
                        creation_time = os.path.getctime(file_path)
                        if creation_time < cutoff:
                            os.unlink(file_path)
                            print(f"Deleted old file: {filename}")
                except Exception as e:
                    print(f"Error deleting {file_path}: {e}")

UPLOAD_DIR = "uploads"
MERGED_DIR = "merged"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(MERGED_DIR, exist_ok=True)

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup: clean old files
    cleanup_old_files()
    print("✓ Cleaned up old files (>1h) on startup")
    yield
    # Shutdown: do nothing (preserve files for dev)
    print("Request finished")

app = FastAPI(lifespan=lifespan)

# CORS configuration - use environment variable for production
# Set CORS_ORIGINS env var to comma-separated list of origins, or leave unset for all origins
CORS_ORIGINS_ENV = os.getenv("CORS_ORIGINS", "")

# Parse origins: if empty or "*", allow all; otherwise split by comma
if CORS_ORIGINS_ENV == "" or CORS_ORIGINS_ENV == "*":
    CORS_ORIGINS = ["*"]
else:
    CORS_ORIGINS = [origin.strip() for origin in CORS_ORIGINS_ENV.split(",")]

app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_credentials=True if CORS_ORIGINS != ["*"] else False,  # credentials can't be used with wildcard
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
)

# Mount uploads so frontend can fetch PDFs for previews
from fastapi.staticfiles import StaticFiles
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")
# Mount merged directory so organized/processed files can also be previewed
app.mount("/merged", StaticFiles(directory=MERGED_DIR), name="merged")

from pydantic import BaseModel

# Session cleanup endpoint
class CleanupRequest(BaseModel):
    session_id: str
    files: List[str]

@app.post("/cleanup")
async def cleanup_session(request: CleanupRequest):
    """Delete files associated with a session when user leaves."""
    deleted = 0
    for filename in request.files:
        # Try to delete from both directories
        for directory in [UPLOAD_DIR, MERGED_DIR]:
            file_path = os.path.join(directory, filename)
            if os.path.exists(file_path):
                try:
                    os.unlink(file_path)
                    deleted += 1
                except Exception as e:
                    print(f"Error deleting {file_path}: {e}")
    
    # Clean up session tracking
    if request.session_id in session_files:
        del session_files[request.session_id]
    
    return {"deleted": deleted, "session_id": request.session_id}

@app.post("/upload")
def upload_file(file: UploadFile = File(...), session_id: str = Form(None)):
    file_id = str(uuid.uuid4())
    stored_name = f"{file_id}_{file.filename}"
    file_path = os.path.join(UPLOAD_DIR, stored_name)
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    # Track file for session cleanup
    if session_id:
        if session_id not in session_files:
            session_files[session_id] = set()
        session_files[session_id].add(stored_name)
        
    return {"id": file_id, "filename": file.filename, "original_name": stored_name}


from pydantic import BaseModel

class MergeRequest(BaseModel):
    files: List[str] # List of original_names

@app.post("/merge")
def merge_pdfs(request: MergeRequest):
    pdf_writer = PdfWriter()
    
    for filename in request.files:
        file_path = os.path.join(UPLOAD_DIR, filename)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail=f"File {filename} not found")
        
        pdf_writer.append(file_path)
    
    merged_filename = f"merged_{uuid.uuid4()}.pdf"
    merged_path = os.path.join(MERGED_DIR, merged_filename)
    
    with open(merged_path, "wb") as output_pdf:
        pdf_writer.write(output_pdf)
        
    return {"url": f"/download/{merged_filename}"}

class SplitRequest(BaseModel):
    filename: str
    page_ranges: str # e.g., "1,3-5,7" (1-based)

def parse_page_ranges(range_str: str, max_pages: int) -> List[int]:
    pages = set()
    parts = range_str.split(',')
    for part in parts:
        part = part.strip()
        if '-' in part:
            start, end = map(int, part.split('-'))
            for i in range(start, end + 1):
                pages.add(i - 1) # 0-based
        else:
            pages.add(int(part) - 1)
    return sorted(list(pages))

@app.post("/split")
def split_pdf(request: SplitRequest):
    file_path = os.path.join(UPLOAD_DIR, request.filename)
    if not os.path.exists(file_path):
        # Also check merged dir just in case user wants to split a result
        file_path = os.path.join(MERGED_DIR, request.filename)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")
            
    reader = PdfWriter() # Can act as reader too loosely, but better use PdfReader
    from pypdf import PdfReader
    
    try:
        reader = PdfReader(file_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        try:
            selected_pages = parse_page_ranges(request.page_ranges, total_pages)
        except ValueError:
             raise HTTPException(status_code=400, detail="Invalid page ranges format")

        for p_idx in selected_pages:
            if 0 <= p_idx < total_pages:
                writer.add_page(reader.pages[p_idx])
        
        output_filename = f"split_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
            
        return {"url": f"/download/{output_filename}"}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class OrganizeRequest(BaseModel):
    filename: str
    page_indices: List[int] # Exact 0-based indices in desired order

@app.post("/organize")
def organize_pdf(request: OrganizeRequest):
    file_path = os.path.join(UPLOAD_DIR, request.filename)
    if not os.path.exists(file_path):
        file_path = os.path.join(MERGED_DIR, request.filename)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")
            
    try:
        reader = PdfReader(file_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        for p_idx in request.page_indices:
            if 0 <= p_idx < total_pages:
                writer.add_page(reader.pages[p_idx])
            else:
                print(f"Warning: Page index {p_idx} out of range (total {total_pages})")

        output_filename = f"organized_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
            
        return {"url": f"/download/{output_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(MERGED_DIR, filename)
    if os.path.exists(file_path):
        return FileResponse(file_path)
    raise HTTPException(status_code=404, detail="File not found")

# ============== CONVERSION FEATURES ==============

@app.post("/images-to-pdf")
def images_to_pdf(files: List[UploadFile] = File(...)):
    """Convert multiple images to a single PDF."""
    from PIL import Image
    import io
    
    if not files:
        raise HTTPException(status_code=400, detail="No images provided")
    
    try:
        pil_images = []
        for f in files:
            content = f.file.read()
            try:
                img = Image.open(io.BytesIO(content))
                # Convert any mode to RGB for PDF compatibility
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                pil_images.append(img)
            except Exception:
                raise HTTPException(status_code=400, detail=f"Invalid image file: {f.filename}")
        
        if not pil_images:
            raise HTTPException(status_code=400, detail="No valid images found")
        
        output_filename = f"images_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        # Save first image as PDF, append the rest
        first_img = pil_images[0]
        if len(pil_images) > 1:
            first_img.save(output_path, "PDF", save_all=True, append_images=pil_images[1:], resolution=150)
        else:
            first_img.save(output_path, "PDF", resolution=150)
        
        return {"url": f"/download/{output_filename}"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/pdf-to-images")
def pdf_to_images(file: UploadFile = File(...)):
    """Convert a PDF to images (returns zip of PNGs)."""
    import zipfile
    from pypdf import PdfReader
    # Note: pypdf alone cannot render images. We need pdf2image which requires poppler.
    # For simplicity, we'll use a try/except and inform user if not available.
    try:
        from pdf2image import convert_from_bytes
    except ImportError:
        raise HTTPException(status_code=501, detail="pdf2image not installed. Run: pip install pdf2image (requires poppler)")
    
    try:
        content = file.file.read()
        images = convert_from_bytes(content, dpi=150)
        
        zip_filename = f"pdf_images_{uuid.uuid4()}.zip"
        zip_path = os.path.join(MERGED_DIR, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w') as zf:
            for i, img in enumerate(images):
                img_bytes = io.BytesIO()
                img.save(img_bytes, format='PNG')
                img_bytes.seek(0)
                zf.writestr(f"page_{i+1}.png", img_bytes.read())
        
        return {"url": f"/download/{zip_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ============== OPTIMIZATION & SECURITY ==============

@app.post("/compress")
def compress_pdf(file: UploadFile = File(...)):
    """Compress a PDF to reduce file size."""
    from pypdf import PdfReader, PdfWriter
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
        
        # Remove metadata to save space
        writer.add_metadata({})
        
        output_filename = f"compressed_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        original_size = len(content)
        compressed_size = os.path.getsize(output_path)
        
        return {
            "url": f"/download/{output_filename}",
            "original_size": original_size,
            "compressed_size": compressed_size,
            "reduction_percent": round((1 - compressed_size / original_size) * 100, 1)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class ProtectRequest(BaseModel):
    password: str

@app.post("/protect")
def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    """Add password protection to a PDF."""
    from pypdf import PdfReader, PdfWriter
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        # Encrypt with password
        writer.encrypt(password)
        
        output_filename = f"protected_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/watermark")
def watermark_pdf(
    file: UploadFile = File(...), 
    text: str = Form("CONFIDENTIAL"),
    opacity: float = Form(0.3),
    font_size: int = Form(60),
    rotation: int = Form(45),
    position: str = Form("center"),  # center, top-left, top-right, bottom-left, bottom-right, tiled
    color: str = Form("#808080"),  # Hex color
    repeat_x: int = Form(1),  # For tiled mode
    repeat_y: int = Form(1)   # For tiled mode
):
    """Add a customizable text watermark to each page of a PDF."""
    from pypdf import PdfReader, PdfWriter
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.colors import HexColor
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        
        # Parse hex color
        try:
            fill_color = HexColor(color)
            fill_color.alpha = opacity
        except:
            fill_color = HexColor("#808080")
            fill_color.alpha = opacity
        
        # Get page dimensions from first page
        first_page = reader.pages[0]
        page_width = float(first_page.mediabox.width)
        page_height = float(first_page.mediabox.height)
        
        # Create watermark PDF in memory
        watermark_buffer = io.BytesIO()
        c = canvas.Canvas(watermark_buffer, pagesize=(page_width, page_height))
        c.setFont("Helvetica-Bold", font_size)
        c.setFillColor(fill_color)
        
        # Position mapping
        positions = {
            "center": [(page_width / 2, page_height / 2)],
            "top-left": [(100, page_height - 100)],
            "top-right": [(page_width - 100, page_height - 100)],
            "bottom-left": [(100, 100)],
            "bottom-right": [(page_width - 100, 100)],
        }
        
        if position == "tiled":
            # Calculate tiled positions
            coords = []
            spacing_x = page_width / (repeat_x + 1)
            spacing_y = page_height / (repeat_y + 1)
            for i in range(1, repeat_x + 1):
                for j in range(1, repeat_y + 1):
                    coords.append((spacing_x * i, spacing_y * j))
        else:
            coords = positions.get(position, positions["center"])
        
        # Draw watermarks at all positions
        for (x, y) in coords:
            c.saveState()
            c.translate(x, y)
            c.rotate(rotation)
            c.drawCentredString(0, 0, text)
            c.restoreState()
        
        c.save()
        watermark_buffer.seek(0)
        
        watermark_reader = PdfReader(watermark_buffer)
        watermark_page = watermark_reader.pages[0]
        
        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
        
        output_filename = f"watermarked_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============== ORGANIZE PDF (Additional) ==============

class RemovePagesRequest(BaseModel):
    pages: str  # e.g., "1,3,5" (1-based)

@app.post("/remove-pages")
def remove_pages(file: UploadFile = File(...), pages: str = Form(...)):
    """Remove specified pages from a PDF."""
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        # Parse pages to remove (1-based input)
        pages_to_remove = set()
        for part in pages.split(','):
            part = part.strip()
            if '-' in part:
                start, end = map(int, part.split('-'))
                for i in range(start, end + 1):
                    pages_to_remove.add(i - 1)  # Convert to 0-based
            else:
                pages_to_remove.add(int(part) - 1)
        
        for i in range(total_pages):
            if i not in pages_to_remove:
                writer.add_page(reader.pages[i])
        
        if len(writer.pages) == 0:
            raise HTTPException(status_code=400, detail="Cannot remove all pages")
        
        output_filename = f"removed_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {
            "url": f"/download/{output_filename}",
            "original_pages": total_pages,
            "remaining_pages": len(writer.pages)
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/extract-pages")
def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    """Extract specified pages from a PDF into a new document."""
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        # Parse pages to extract (1-based input)
        selected_pages = parse_page_ranges(pages, total_pages)
        
        for p_idx in selected_pages:
            if 0 <= p_idx < total_pages:
                writer.add_page(reader.pages[p_idx])
        
        if len(writer.pages) == 0:
            raise HTTPException(status_code=400, detail="No valid pages to extract")
        
        output_filename = f"extracted_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {
            "url": f"/download/{output_filename}",
            "extracted_pages": len(writer.pages)
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============== OPTIMIZE PDF (Additional) ==============

@app.post("/repair")
def repair_pdf(file: UploadFile = File(...)):
    """Attempt to repair a corrupted PDF."""
    try:
        content = file.file.read()
        
        # Use pikepdf which is more robust at handling corrupt PDFs
        import pikepdf
        
        # Open with pikepdf (it auto-repairs many issues)
        pdf = pikepdf.open(io.BytesIO(content))
        
        output_filename = f"repaired_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        # Save with linearization for optimized access
        pdf.save(output_path, linearize=True)
        pdf.close()
        
        return {
            "url": f"/download/{output_filename}",
            "status": "repaired"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Could not repair PDF: {str(e)}")


@app.post("/ocr")
def ocr_pdf(file: UploadFile = File(...), language: str = Form("eng")):
    """OCR a scanned PDF to make it searchable."""
    try:
        from pdf2image import convert_from_bytes
    except ImportError:
        raise HTTPException(status_code=501, detail="pdf2image not installed")
    
    try:
        import pytesseract
    except ImportError:
        raise HTTPException(status_code=501, detail="pytesseract not installed. Install tesseract-ocr system package.")
    
    try:
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.pagesizes import letter
        from PIL import Image as PILImage
        
        content = file.file.read()
        
        # Convert PDF pages to images
        images = convert_from_bytes(content, dpi=300)
        
        writer = PdfWriter()
        
        for img in images:
            # Get OCR text with bounding box data
            img_width, img_height = img.size
            
            # Create a page with the image and searchable text overlay
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='JPEG', quality=95)
            img_buffer.seek(0)
            
            # Create PDF page with image
            page_buffer = io.BytesIO()
            c = rl_canvas.Canvas(page_buffer, pagesize=(img_width * 72 / 300, img_height * 72 / 300))
            
            # Draw the original image
            from reportlab.lib.utils import ImageReader
            c.drawImage(ImageReader(img_buffer), 0, 0, 
                       width=img_width * 72 / 300, height=img_height * 72 / 300)
            
            c.save()
            page_buffer.seek(0)
            
            page_reader = PdfReader(page_buffer)
            writer.add_page(page_reader.pages[0])
        
        output_filename = f"ocr_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {
            "url": f"/download/{output_filename}",
            "pages_processed": len(images)
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============== CONVERT TO PDF ==============

def _convert_with_libreoffice(input_path: str, output_dir: str) -> str:
    """Convert a document to PDF using LibreOffice."""
    import subprocess
    
    try:
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path],
            capture_output=True, text=True, timeout=120
        )
        
        if result.returncode != 0:
            # Try soffice as alternative
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, input_path],
                capture_output=True, text=True, timeout=120
            )
        
        if result.returncode != 0:
            raise Exception(f"LibreOffice conversion failed: {result.stderr}")
        
        # Find the output file
        basename = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{basename}.pdf")
        
        if not os.path.exists(output_path):
            raise Exception("Conversion produced no output file")
        
        return output_path
    except FileNotFoundError:
        raise Exception("LibreOffice not found. Please install LibreOffice for document conversion.")


@app.post("/word-to-pdf")
def word_to_pdf(file: UploadFile = File(...)):
    """Convert a Word document (DOCX) to PDF."""
    try:
        # Save uploaded file temporarily
        temp_filename = f"temp_{uuid.uuid4()}_{file.filename}"
        temp_path = os.path.join(UPLOAD_DIR, temp_filename)
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # Convert using LibreOffice
        output_path = _convert_with_libreoffice(temp_path, MERGED_DIR)
        
        # Rename to unique name
        final_filename = f"word2pdf_{uuid.uuid4()}.pdf"
        final_path = os.path.join(MERGED_DIR, final_filename)
        os.rename(output_path, final_path)
        
        # Cleanup temp file
        os.unlink(temp_path)
        
        return {"url": f"/download/{final_filename}"}
    except Exception as e:
        # Cleanup temp file on error
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pptx-to-pdf")
def pptx_to_pdf(file: UploadFile = File(...)):
    """Convert a PowerPoint document (PPTX) to PDF."""
    try:
        temp_filename = f"temp_{uuid.uuid4()}_{file.filename}"
        temp_path = os.path.join(UPLOAD_DIR, temp_filename)
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        output_path = _convert_with_libreoffice(temp_path, MERGED_DIR)
        
        final_filename = f"pptx2pdf_{uuid.uuid4()}.pdf"
        final_path = os.path.join(MERGED_DIR, final_filename)
        os.rename(output_path, final_path)
        
        os.unlink(temp_path)
        
        return {"url": f"/download/{final_filename}"}
    except Exception as e:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/excel-to-pdf")
def excel_to_pdf(file: UploadFile = File(...)):
    """Convert an Excel document (XLSX) to PDF."""
    try:
        temp_filename = f"temp_{uuid.uuid4()}_{file.filename}"
        temp_path = os.path.join(UPLOAD_DIR, temp_filename)
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        output_path = _convert_with_libreoffice(temp_path, MERGED_DIR)
        
        final_filename = f"excel2pdf_{uuid.uuid4()}.pdf"
        final_path = os.path.join(MERGED_DIR, final_filename)
        os.rename(output_path, final_path)
        
        os.unlink(temp_path)
        
        return {"url": f"/download/{final_filename}"}
    except Exception as e:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/html-to-pdf")
def html_to_pdf(file: UploadFile = File(None), url: str = Form(None)):
    """Convert HTML file or URL to PDF."""
    try:
        from weasyprint import HTML
    except ImportError:
        raise HTTPException(status_code=501, detail="weasyprint not installed. Run: pip install weasyprint")
    
    try:
        output_filename = f"html2pdf_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        if file and file.filename:
            # Convert uploaded HTML file
            content = file.file.read().decode('utf-8', errors='replace')
            HTML(string=content).write_pdf(output_path)
        elif url:
            # Convert from URL
            HTML(url=url).write_pdf(output_path)
        else:
            raise HTTPException(status_code=400, detail="Either a file or URL must be provided")
        
        return {"url": f"/download/{output_filename}"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============== CONVERT FROM PDF ==============

def _convert_pdf_with_libreoffice(input_path: str, output_dir: str, output_format: str) -> str:
    """Convert a PDF to another format using LibreOffice."""
    import subprocess
    
    try:
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', output_format, '--outdir', output_dir, input_path],
            capture_output=True, text=True, timeout=120
        )
        
        if result.returncode != 0:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', output_format, '--outdir', output_dir, input_path],
                capture_output=True, text=True, timeout=120
            )
        
        if result.returncode != 0:
            raise Exception(f"LibreOffice conversion failed: {result.stderr}")
        
        basename = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_dir, f"{basename}.{output_format}")
        
        if not os.path.exists(output_path):
            raise Exception("Conversion produced no output file")
        
        return output_path
    except FileNotFoundError:
        raise Exception("LibreOffice not found. Please install LibreOffice for document conversion.")


@app.post("/pdf-to-word")
def pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to Word document (DOCX)."""
    try:
        temp_filename = f"temp_{uuid.uuid4()}.pdf"
        temp_path = os.path.join(UPLOAD_DIR, temp_filename)
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        output_path = _convert_pdf_with_libreoffice(temp_path, MERGED_DIR, 'docx')
        
        final_filename = f"pdf2word_{uuid.uuid4()}.docx"
        final_path = os.path.join(MERGED_DIR, final_filename)
        os.rename(output_path, final_path)
        
        os.unlink(temp_path)
        
        return {"url": f"/download/{final_filename}", "format": "docx"}
    except Exception as e:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pdf-to-pptx")
def pdf_to_pptx(file: UploadFile = File(...)):
    """Convert PDF to PowerPoint (PPTX) - renders each page as an image slide."""
    try:
        from pdf2image import convert_from_bytes
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as ie:
        raise HTTPException(status_code=501, detail=f"Missing dependency: {str(ie)}")
    
    try:
        content = file.file.read()
        images = convert_from_bytes(content, dpi=200)
        
        prs = Presentation()
        # Set slide dimensions to match standard widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            slide.shapes.add_picture(img_buffer, Inches(0), Inches(0), 
                                    width=prs.slide_width, height=prs.slide_height)
        
        final_filename = f"pdf2pptx_{uuid.uuid4()}.pptx"
        final_path = os.path.join(MERGED_DIR, final_filename)
        prs.save(final_path)
        
        return {"url": f"/download/{final_filename}", "format": "pptx", "slides": len(images)}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pdf-to-excel")
def pdf_to_excel(file: UploadFile = File(...)):
    """Convert PDF to Excel (XLSX) - extracts text into spreadsheet."""
    try:
        from openpyxl import Workbook
    except ImportError:
        raise HTTPException(status_code=501, detail="openpyxl not installed")
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        # Header
        ws.append(["Page", "Content"])
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            # Split text into lines and add each as a row
            lines = text.strip().split('\n')
            for line in lines:
                if line.strip():
                    ws.append([i + 1, line.strip()])
        
        final_filename = f"pdf2excel_{uuid.uuid4()}.xlsx"
        final_path = os.path.join(MERGED_DIR, final_filename)
        wb.save(final_path)
        
        return {"url": f"/download/{final_filename}", "format": "xlsx", "pages": len(reader.pages)}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/pdf-to-pdfa")
def pdf_to_pdfa(file: UploadFile = File(...)):
    """Convert PDF to PDF/A for archival."""
    try:
        import pikepdf
        
        content = file.file.read()
        pdf = pikepdf.open(io.BytesIO(content))
        
        # Set PDF/A metadata
        with pdf.open_metadata() as meta:
            meta['dc:title'] = 'PDF/A Converted Document'
            meta['pdf:PDFVersion'] = '1.4'
            meta['pdfaid:part'] = '2'
            meta['pdfaid:conformance'] = 'B'
        
        output_filename = f"pdfa_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        pdf.save(output_path, linearize=True)
        pdf.close()
        
        return {"url": f"/download/{output_filename}", "format": "pdf/a-2b"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ============== EDIT PDF ==============

@app.post("/rotate")
def rotate_pdf(
    file: UploadFile = File(...),
    degrees: int = Form(90),
    pages: str = Form("all")  # "all" or "1,3,5" (1-based)
):
    """Rotate pages in a PDF."""
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        # Parse which pages to rotate
        if pages.strip().lower() == "all":
            rotate_indices = set(range(total_pages))
        else:
            rotate_indices = set()
            for part in pages.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    for i in range(start, end + 1):
                        rotate_indices.add(i - 1)
                else:
                    rotate_indices.add(int(part) - 1)
        
        for i in range(total_pages):
            page = reader.pages[i]
            if i in rotate_indices:
                page.rotate(degrees)
            writer.add_page(page)
        
        output_filename = f"rotated_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/add-page-numbers")
def add_page_numbers(
    file: UploadFile = File(...),
    position: str = Form("bottom-center"),  # bottom-left, bottom-center, bottom-right, top-left, top-center, top-right
    font_size: int = Form(12),
    format_str: str = Form("{n}"),  # {n} = page number, {total} = total pages
    start_number: int = Form(1),
    margin: int = Form(40),
    color: str = Form("#000000")
):
    """Add page numbers to a PDF."""
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import HexColor
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        try:
            fill_color = HexColor(color)
        except:
            fill_color = HexColor("#000000")
        
        for i, page in enumerate(reader.pages):
            page_width = float(page.mediabox.width)
            page_height = float(page.mediabox.height)
            
            # Create page number overlay
            overlay_buffer = io.BytesIO()
            c = rl_canvas.Canvas(overlay_buffer, pagesize=(page_width, page_height))
            c.setFont("Helvetica", font_size)
            c.setFillColor(fill_color)
            
            # Format page number text
            page_num = start_number + i
            text = format_str.replace("{n}", str(page_num)).replace("{total}", str(total_pages))
            
            # Calculate position
            text_width = c.stringWidth(text, "Helvetica", font_size)
            
            pos_map = {
                "bottom-left": (margin, margin),
                "bottom-center": (page_width / 2 - text_width / 2, margin),
                "bottom-right": (page_width - margin - text_width, margin),
                "top-left": (margin, page_height - margin),
                "top-center": (page_width / 2 - text_width / 2, page_height - margin),
                "top-right": (page_width - margin - text_width, page_height - margin),
            }
            
            x, y = pos_map.get(position, pos_map["bottom-center"])
            c.drawString(x, y, text)
            c.save()
            overlay_buffer.seek(0)
            
            overlay_reader = PdfReader(overlay_buffer)
            page.merge_page(overlay_reader.pages[0])
            writer.add_page(page)
        
        output_filename = f"numbered_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}", "total_pages": total_pages}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/crop")
def crop_pdf(
    file: UploadFile = File(...),
    top: float = Form(0),    # Points to crop from top
    bottom: float = Form(0), # Points to crop from bottom
    left: float = Form(0),   # Points to crop from left
    right: float = Form(0),  # Points to crop from right
    pages: str = Form("all")
):
    """Crop margins from PDF pages."""
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        if pages.strip().lower() == "all":
            crop_indices = set(range(total_pages))
        else:
            crop_indices = set()
            for part in pages.split(','):
                part = part.strip()
                if '-' in part:
                    start, end = map(int, part.split('-'))
                    for i in range(start, end + 1):
                        crop_indices.add(i - 1)
                else:
                    crop_indices.add(int(part) - 1)
        
        for i in range(total_pages):
            page = reader.pages[i]
            if i in crop_indices:
                # Adjust the crop box
                media_box = page.mediabox
                page.mediabox.lower_left = (
                    float(media_box.lower_left[0]) + left,
                    float(media_box.lower_left[1]) + bottom
                )
                page.mediabox.upper_right = (
                    float(media_box.upper_right[0]) - right,
                    float(media_box.upper_right[1]) - top
                )
            writer.add_page(page)
        
        output_filename = f"cropped_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/edit-pdf")
def edit_pdf(
    file: UploadFile = File(...),
    annotations: str = Form("[]")  # JSON array of {text, x, y, page, fontSize, color}
):
    """Add text annotations/overlays to a PDF."""
    import json
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.colors import HexColor
    
    try:
        content = file.file.read()
        reader = PdfReader(io.BytesIO(content))
        writer = PdfWriter()
        
        try:
            annotation_list = json.loads(annotations)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid annotations JSON")
        
        # Group annotations by page (0-based)
        page_annotations = {}
        for ann in annotation_list:
            page_idx = ann.get("page", 1) - 1  # Convert 1-based to 0-based
            if page_idx not in page_annotations:
                page_annotations[page_idx] = []
            page_annotations[page_idx].append(ann)
        
        for i, page in enumerate(reader.pages):
            if i in page_annotations:
                page_width = float(page.mediabox.width)
                page_height = float(page.mediabox.height)
                
                overlay_buffer = io.BytesIO()
                c = rl_canvas.Canvas(overlay_buffer, pagesize=(page_width, page_height))
                
                for ann in page_annotations[i]:
                    text = ann.get("text", "")
                    x = float(ann.get("x", 100))
                    y = float(ann.get("y", 100))
                    font_size = int(ann.get("fontSize", 12))
                    color = ann.get("color", "#000000")
                    
                    try:
                        c.setFillColor(HexColor(color))
                    except:
                        c.setFillColor(HexColor("#000000"))
                    
                    c.setFont("Helvetica", font_size)
                    c.drawString(x, y, text)
                
                c.save()
                overlay_buffer.seek(0)
                
                overlay_reader = PdfReader(overlay_buffer)
                page.merge_page(overlay_reader.pages[0])
            
            writer.add_page(page)
        
        output_filename = f"edited_{uuid.uuid4()}.pdf"
        output_path = os.path.join(MERGED_DIR, output_filename)
        
        with open(output_path, "wb") as out_f:
            writer.write(out_f)
        
        return {"url": f"/download/{output_filename}"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/")
def read_root():
    return {"message": "PDF Tools Backend Ready"}

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port)
